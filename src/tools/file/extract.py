import os
import json
from pptx import Presentation
import kreuzberg
from src.security.path_safety import resolve_safe_path
from src.utils.exceptions import SecurityException


def _extract_pdf(file_path: str) -> dict:
    result = kreuzberg.extract_file_sync(file_path)
    pages_text = []
    if result.pages:
        for page in result.pages:
            pages_text.append(page.content if hasattr(page, 'content') else str(page))
    content = result.content or "\n".join(pages_text)
    tables_data = []
    if result.tables:
        for table in result.tables:
            if hasattr(table, 'rows'):
                rows = [[str(cell) for cell in row] for row in table.rows]
                tables_data.append(rows)
    return {
        "content": content,
        "page_count": len(result.pages) if result.pages else 0,
        "table_count": len(tables_data),
        "tables": tables_data,
    }


def _extract_ppt(file_path: str) -> dict:
    prs = Presentation(file_path)
    slides_data = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_texts.append(row_text)
        slides_data.append({"slide": slide_num, "text": "\n".join(slide_texts)})
    full_text = "\n\n".join(
        f"--- 第 {s['slide']} 页 ---\n{s['text']}" for s in slides_data if s['text']
    )
    return {
        "content": full_text,
        "slide_count": len(prs.slides),
        "slides": slides_data,
    }


def read_pdf(path: str, filename: str) -> str:
    if not filename.lower().endswith(".pdf"):
        return json.dumps(
            {
                "success": False,
                "summary": "仅支持 .pdf",
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )
    path = os.path.abspath(path)
    file_path = os.path.join(path, filename)
    try:
        file_path = resolve_safe_path(file_path)
    except SecurityException as e:
        return json.dumps(
            {
                "success": False,
                "summary": str(e),
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )
    if not os.path.isfile(file_path):
        return json.dumps(
            {
                "success": False,
                "summary": f"文件不存在: {file_path}",
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )
    try:
        extracted = _extract_pdf(file_path)
        return json.dumps(
            {
                "success": True,
                "summary": f"PDF 提取完成: {extracted['page_count']} 页 {extracted['table_count']} 表",
                "content": extracted["content"],
                "data": {"page_count": extracted["page_count"], "table_count": extracted["table_count"], "tables": extracted["tables"]}
            },
            ensure_ascii=False
        )
    except Exception as e:
        return json.dumps(
            {
                "success": False,
                "summary": str(e),
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )


def read_ppt(path: str, filename: str) -> str:
    lower = filename.lower()
    if not (lower.endswith(".pptx") or lower.endswith(".ppt")):
        return json.dumps(
            {
                "success": False,
                "summary": "仅支持 .pptx",
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )
    if lower.endswith(".ppt"):
        return json.dumps(
            {
                "success": False,
                "summary": "暂不支持旧版 .ppt，请转换为 .pptx",
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )
    path = os.path.abspath(path)
    file_path = os.path.join(path, filename)
    try:
        file_path = resolve_safe_path(file_path)
    except SecurityException as e:
        return json.dumps(
            {
                "success": False,
                "summary": str(e),
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )
    if not os.path.isfile(file_path):
        return json.dumps(
            {
                "success": False,
                "summary": f"文件不存在: {file_path}",
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )
    try:
        extracted = _extract_ppt(file_path)
        return json.dumps(
            {
                "success": True,
                "summary": f"PPT 提取完成: {extracted['slide_count']} 张幻灯片",
                "content": extracted["content"],
                "data": {"slide_count": extracted["slide_count"], "slides": extracted["slides"]}
            },
            ensure_ascii=False
        )
    except Exception as e:
        return json.dumps(
            {
                "success": False,
                "summary": str(e),
                "content": None,
                "data": None
            },
            ensure_ascii=False
        )


READ_PDF_SCHEMA = {
    "type": "function",
    "function": {
        "name": "read_pdf",
        "description": "提取PDF文件的全部文本内容和表格",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "父目录绝对路径"},
                "filename": {"type": "string", "description": ".pdf 文件名"},
            },
            "required": ["path", "filename"],
        },
    },
}

READ_PPT_SCHEMA = {
    "type": "function",
    "function": {
        "name": "read_ppt",
        "description": "提取PPT/PPTX文件的全部幻灯片文本和表格内容（仅支持.pptx，不支持旧版.ppt）",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "父目录绝对路径"},
                "filename": {"type": "string", "description": ".pptx 文件名"},
            },
            "required": ["path", "filename"],
        },
    },
}
